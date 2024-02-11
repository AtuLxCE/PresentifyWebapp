from fastapi import FastAPI, UploadFile, File
import pdftitle, pdfplumber
from bs4 import BeautifulSoup
from classobjects import PDF, PresentationData, TextSummaryRequest, TextSummaryResponse
from pdftools import read_pdf, read_pdf_from_url, clean_text
from gemini import gemini_summarize
import requests
import pandas as pd
import difflib
from fastapi import HTTPException
from presentify_model import summarize
from pptx import Presentation
from pptx.util import Inches
from pptxtools import *



pdf = PDF(textdata=None)
presentation = PresentationData(title=None, author=None,
                             abstract=None, introduction=None, literature_review=None, 
                             methodology=None, results=None, conclusions=None)




app = FastAPI()

MAX_FILE_SIZE = 1024 * 1024 * 12  # 12MB 
@app.post('/extract-text')
async def extract_texts(file: UploadFile = File(...)):
    if file.size > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File is too large")
    pdf_bytes = await file.read()

    # Optionally save the PDF temporarily
    with open('temp_pdf.pdf', 'wb') as f:
        f.write(pdf_bytes)

    pdf.textdata = clean_text(read_pdf('temp_pdf.pdf'))  # Replace with your program's function
    pdf_title = pdftitle.get_title_from_file('temp_pdf.pdf')
    author = pdfplumber.open('temp_pdf.pdf').metadata['Author']
    presentation.author = author
    presentation.title = pdf_title
    return {'title':presentation.title,'textdata':pdf.textdata}


@app.post("/get_data_from_url")
async def get_data_fromI_url(arxiv_url: str):
    response = requests.get(arxiv_url)
    response = response.content
    soup = BeautifulSoup(response, 'html.parser')

    presentation.title =  soup.find('h1', class_='title mathjax').text
    presentation.author = soup.find('div', class_='authors').text
    pdf_link = arxiv_url.replace('abs', 'pdf')
    pdf.textdata = clean_text(read_pdf_from_url(pdf_link))
    presentation.title = presentation.title.replace('Title:','')
    presentation.author = presentation.author.replace('Authors:','')
    return{'title':presentation.title,'author':presentation.author,'textdata':pdf.textdata}
    
   
@app.get('/get-section-data')
def get_section_data():
    text = pdf.textdata
    gemini_data = PresentationData()
    try:
        gemini_data = gemini_summarize(text)
        presentation.introduction = gemini_data.introduction
        presentation.literature_review = gemini_data.literature_review
        presentation.methodology = gemini_data.methodology
        presentation.results = gemini_data.results
        presentation.conclusions = gemini_data.conclusions
    except:
        return {'error':'couldnt extract data'}
        
    return {'title': presentation.title, 'author': presentation.author,
            'introduction':gemini_data.introduction,
            'literature_review':gemini_data.literature_review,
            'methodology':gemini_data.methodology,
            'results':gemini_data.results,
            'conclusions':gemini_data.conclusions
            }
    
data = PresentationData(title=None, author=None,
                             abstract=None, introduction=None, literature_review=None, 
                             methodology=None, datas=None, conclusions=None)


title_list = ['Introduction','Literature Review','Methodology','Results','Conclusion']
@app.get('/generate_slide')
def predict():
    data = presentation
    data_dict = {'Introduction':data.introduction,
             'Literature Review':data.literature_review, 
             'Methodology':data.methodology, 
             'Results':data.results, 
             'Conclusion':data.conclusions}
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    #addding page 1
    slide1 = add_slide(prs, title_slide_layout)
    slide1.placeholders[0].text = data.title
    slide1.placeholders[1].text= data.author

    #font style for title
    customizer_placeholder(slide1,0,'Arial',44,True)

    #font style for author
    customizer_placeholder(slide1,1,'Arial',22,True)

    #bullet slides
    for title in title_list:
        slide2 = add_slide(prs, bullet_slide_layout)
        slide2.placeholders[0].text = title

        content = data_dict[title]
        sentences = split_sentences(content)
        print(sentences)

        slide2.shapes.placeholders[1].text_frame.text = sentences[0]
        customizer_placeholder(slide2,1,'Arial',28,False)

        for i in range(1,len(sentences)):
            p = slide2.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = sentences[i]
            p.level = 0
            #font style for bullet
            customizer_placeholder_paragraphs(slide2,1,i,'Arial',28,False)
    #font style for title
    customizer_placeholder(slide2,0,'Arial',36,True)
    prs.save(f'slides/{data.title}.pptx')
    print('slide successfully created')
    return {'title': data.title, 'author': data.author,
            'introduction':data.introduction,
            'literature_review':data.literature_review,
            'methodology':data.methodology,
            'results':data.results,
            'conclusions':data.conclusions,
            }
